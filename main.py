from fastapi import FastAPI, Request, Depends, HTTPException, Cookie
from fastapi.responses import HTMLResponse, JSONResponse, RedirectResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from fastapi import Body, Depends
from chainlit.utils import mount_chainlit
from chainlit.server import app as chainlit_app
from chainlit.auth import get_current_user
from pathlib import Path
from sqlalchemy import create_engine, text
from sqlalchemy.orm import sessionmaker
from QAgent.config.config_manager import config
import json
import jwt
from typing import Optional
import traceback
from uuid import UUID
import pandas as pd
import numpy as np
from datetime import date, timedelta, datetime

# Para la presentación de gráficos 
# asegurarse de instalar pip install plotly[kaleido]
from pptx import Presentation
from pptx.util import Inches
import plotly.io as pio
import uuid
from dotenv import load_dotenv
load_dotenv() 

# Crear la aplicación FastAPI
app = FastAPI()


# Configurar directorios de archivos estáticos
PUBLIC_DIR = Path("public")
STORAGE_DIR = PUBLIC_DIR / "storage"

# Asegurarse de que los directorios existan
PUBLIC_DIR.mkdir(exist_ok=True)
STORAGE_DIR.mkdir(exist_ok=True)

# Montar los directorios de archivos estáticos
app.mount("/public", StaticFiles(directory="public"), name="public")
app.mount("/public/storage", StaticFiles(directory="public/storage"), name="storage")

# Configurar templates
templates = Jinja2Templates(directory="templates")

# Configuración de SQLAlchemy
db_uri = config.get('CHAINLIT_DB_URI')
# print("\n🔧 Configuración de base de datos:")
# print(f"URI de conexión: {db_uri}")

# Si la URI usa asyncpg, cambiarla a psycopg2
if 'postgresql+asyncpg' in db_uri:
    print("⚠️  Detectada URI para asyncpg, cambiando a psycopg2...")
    db_uri = db_uri.replace('postgresql+asyncpg', 'postgresql+psycopg2')
    # print(f"Nueva URI: {db_uri}")

# Crear el engine con echo=True para ver las consultas SQL
engine = create_engine(db_uri, echo=True)
Session = sessionmaker(bind=engine)

# Función auxiliar para serializar UUIDs
def serialize_uuid(obj):
    if isinstance(obj, UUID):
        return str(obj)
    raise TypeError(f'Object of type {obj.__class__.__name__} is not JSON serializable')

def formatea_duracion(segundos):
    if segundos < 60:
        return f"{segundos:.2f} Seg."
    else:
        minutos = int(segundos // 60)
        segundos_restantes = int(segundos % 60)
        return f"{minutos}:{segundos_restantes:02d} Min."

# Función para verificar rol de administrador
def verify_admin_role(user):
    """Verifica que el usuario tenga rol de administrador"""
    if not user:
        raise HTTPException(status_code=401, detail="No autenticado")

    user_metadata = user.get("metadata", {})
    user_role = user_metadata.get("role", "user").lower()

    if user_role != "admin":
        raise HTTPException(status_code=403, detail="Acceso denegado: Se requiere rol de administrador")


# Función para obtener el usuario actual usando la cookie de sesión de Chainlit
@app.get("/api/user")
def get_user_from_cookie(access_token: Optional[str] = Cookie(None)):
    if not access_token:
        print("❌ No se encontró token de acceso")
        raise HTTPException(status_code=401, detail="No autenticado")
    
    try:
        print("\n🔍 Decodificando token JWT...")
        # Decodificar el token sin verificar la firma
        decoded_token = jwt.decode(access_token, options={"verify_signature": False})
        print(f"✅ Contenido del token: {json.dumps(decoded_token, indent=2)}")
        
        # Extraer el identificador del usuario del token
        user_identifier = decoded_token.get('identifier')
        user_metadata = decoded_token.get('metadata', {})
        
        print(f"👤 Usuario identificado: {user_identifier}")
        print(f"📋 Metadata: {json.dumps(user_metadata, indent=2)}")
        
        return {
            "identifier": user_identifier,
            "metadata": user_metadata
        }
        
    except Exception as e:
        print(f"❌ Error decodificando token: {str(e)}")
        raise HTTPException(status_code=401, detail="Error de autenticación")

# Endpoint para obtener gráficos pineados
@app.get("/api/pinned-graphs")
def get_pinned_graphs(user = Depends(get_user_from_cookie)):
    print("\n📊 [get_pinned_graphs] Iniciando búsqueda de gráficos...")
    print(f"👤 Usuario: {user['identifier']}")
    
    try:
        print("🔄 Creando sesión de base de datos...")
        session = Session()
        
        print("🔍 Preparando consulta SQL...")
        # Primero, vamos a buscar el thread correcto
        thread_query = text("""
            SELECT id, "userIdentifier", "createdAt"
            FROM threads 
            WHERE "userIdentifier" = :user_identifier
        """)
        
        print("🔍 Buscando threads del usuario...")
        threads = session.execute(
            thread_query,
            {"user_identifier": user["identifier"]}
        ).fetchall()
        
        print(f"📋 Threads encontrados: {len(threads)}")
        for thread in threads:
            print(f"   - Thread ID: {thread.id}, User: {thread.userIdentifier}, Created: {thread.createdAt}")
        
        print("\n🔍 Buscando elementos pineados...")
        query = text("""
            SELECT 
            e.id as element_id,
            e."threadId",
            e.type,
            e.url,
            e.name,
            e.display,
            e."objectKey",
            e.mime,
            e.props,
            e.pin,
            e.pin_title,
            e.created_at,
            e.pin_update,            -- NUEVO
            e.pin_update_at,         -- NUEVO
            t.id as thread_id,
            t."userIdentifier",
            t."createdAt"
            FROM elements e
            JOIN threads t ON e."threadId" = t.id
            WHERE e.type = 'plotly' 
            AND e.pin = 1
            AND t."userIdentifier" = :user_identifier
            ORDER BY t."createdAt" DESC

        """)
        
        print(f"🔄 Ejecutando consulta para usuario: {user['identifier']}")
        result = session.execute(
            query,
            {"user_identifier": user["identifier"]}
        )
        
        print("✅ Consulta ejecutada, procesando resultados...")
        graphs = []
        for row in result:
            try:
                print(f"\n📄 Procesando gráfico:")
                print(f"   ID: {row.element_id}")
                print(f"   Thread ID: {row.thread_id}")
                print(f"   Object Key: {row.objectKey}")
                
                # Intentar usar objectKey si está disponible
                if row.objectKey:
                    print("📂 Usando objectKey para la ruta...")
                    json_path = PUBLIC_DIR / "storage" / row.objectKey
                else:
                    print("📂 Construyendo ruta con IDs...")
                    thread_dir = str(row.thread_id)
                    file_dir = str(row.element_id)
                    json_path = PUBLIC_DIR / "storage" / thread_dir / file_dir / "chart.json"
                
                print(f"📂 Buscando archivo en: {json_path}")
                if json_path.exists():
                    print(f"✅ Archivo encontrado")
                    with open(json_path) as f:
                        figure_data = json.load(f)
                    print("✅ Datos del gráfico cargados")
                else:
                    print(f"⚠️ Archivo no encontrado")
                    # Intentar buscar en una ubicación alternativa
                    alt_path = PUBLIC_DIR / "storage" / "e9c7c532-2b6d-4308-9c36-ebeda53fb22d" / str(row.element_id) / "chart.json"
                    print(f"🔄 Intentando ubicación alternativa: {alt_path}")
                    if alt_path.exists():
                        print(f"✅ Archivo encontrado en ubicación alternativa")
                        with open(alt_path) as f:
                            figure_data = json.load(f)
                        print("✅ Datos del gráfico cargados")
                    else:
                        print(f"❌ Archivo no encontrado en ninguna ubicación")
                        figure_data = None
                
               # Elegir fecha a mostrar: si pin_update=1 y hay pin_update_at => úsala, si no created_at
                display_date = None
                try:
                    if getattr(row, "pin_update", None) == 1 and getattr(row, "pin_update_at", None):
                        display_date = row.pin_update_at  # ya viene varchar 'YYYY-MM-DD'
                    elif row.created_at:
                        display_date = row.created_at.strftime("%d-%m-%Y")
                except Exception:
                    display_date = row.created_at.strftime("%d-%m-%Y") if row.created_at else None

                graph_data = {
                    "id": str(row.element_id),
                    "threadId": str(row.thread_id),
                    "type": row.type,
                    "url": row.url,
                    "name": row.name,
                    "display": row.display,
                    "figure": figure_data,
                    "pin_title": row.pin_title,
                    "created_at": row.created_at.strftime("%d-%m-%Y") if row.created_at else None,
                    "objectKey": row.objectKey,                               # NUEVO
                    "pin_update": int(row.pin_update) if row.pin_update is not None else 0,  # NUEVO
                    "pin_update_at": row.pin_update_at,                       # NUEVO (varchar)
                    "display_date": display_date                              # NUEVO (lo usaremos en el header)
                }
                graphs.append(graph_data)
                print("✅ Gráfico agregado a la lista")
            except Exception as e:
                print(f"❌ Error procesando gráfico {row.element_id}: {str(e)}")
                continue
        
        print(f"\n✅ Proceso completado. Total de gráficos: {len(graphs)}")
        return JSONResponse(
            content={"graphs": graphs},
            status_code=200
        )
    except Exception as e:
        print(f"❌ Error general: {str(e)}")
        print(f"Detalles del error: {traceback.format_exc()}")
        return JSONResponse(
            status_code=500,
            content={"error": str(e)}
        )
    finally:
        if 'session' in locals():
            print("🔄 Cerrando sesión de base de datos...")
            session.close()
            print("✅ Sesión cerrada")


@app.post("/api/delete-graph")
def delete_graph(payload: dict = Body(...), user=Depends(get_user_from_cookie)):
    try:
        graph_id = payload.get("id")
        print(f"🧹 Solicitud para eliminar gráfico ID: {graph_id}")
        
        # Crear sesión
        session = Session()

        # Ejecutar UPDATE en la tabla elements
        update_query = text("""
            UPDATE elements
            SET pin = 0
            WHERE id = :graph_id
        """)
        result = session.execute(update_query, {"graph_id": graph_id})
        session.commit()

        print("✅ Gráfico despineado con éxito")
        return JSONResponse(content={"status": "ok", "message": "Gráfico despineado"}, status_code=200)
    
    except Exception as e:
        print(f"❌ Error en delete_graph: {e}")
        return JSONResponse(content={"error": str(e)}, status_code=500)
    
    finally:
        if 'session' in locals():
            session.close()
            print("✅ Sesión cerrada")
            
            
@app.post("/api/update-graph")
def update_graph(payload: dict = Body(...), user=Depends(get_user_from_cookie)):
    try:
        graph_id = payload.get("id")
        raw_title = payload.get("title", "")
        pin_title = raw_title.strip()

        print(f"🧹 Solicitud para actualizar gráfico ID: {graph_id}")
        print(f"📝 Nuevo título recibido (raw): '{raw_title}'")

        # Validación: título no puede estar vacío después de hacer strip
        if not pin_title:
            print("❌ Título vacío o solo espacios, no se actualiza")
            return JSONResponse(
                content={"error": "El título no puede estar vacío."},
                status_code=400
            )

        # Crear sesión
        session = Session()

        # Ejecutar UPDATE
        update_query = text("""
            UPDATE elements
            SET pin_title = :pin_title
            WHERE id = :graph_id
        """)
        session.execute(update_query, {"pin_title": pin_title, "graph_id": graph_id})
        session.commit()

        print("✅ Gráfico actualizado con éxito")
        return JSONResponse(
            content={"status": "ok", "message": "Título actualizado correctamente"},
            status_code=200
        )

    except Exception as e:
        print(f"❌ Error en update_graph: {e}")
        return JSONResponse(content={"error": str(e)}, status_code=500)

    finally:
        if 'session' in locals():
            session.close()
            print("✅ Sesión cerrada")

from pptx import Presentation
from pptx.util import Inches
import plotly.io as pio
import uuid

@app.post("/api/create_powerpoint")
def create_powerpoint(payload: dict = Body(...), user=Depends(get_user_from_cookie)):
    try:
        graph_ids = payload.get("graph_ids", [])
        if not graph_ids:
            return JSONResponse(content={"error": "No se recibieron IDs de gráficos"}, status_code=400)

        user_identifier = user["identifier"]
        print(f"👤 Usuario identificado: {user_identifier}" )
        
        
        session = Session()
        try:
            user_query = text("""
                SELECT id FROM users
                WHERE "identifier" = :user_identifier
                LIMIT 1
            """)
            result = session.execute(user_query, {"user_identifier": user_identifier}).fetchone()

            if not result:
                return JSONResponse(content={"error": "Usuario no encontrado"}, status_code=404)

            user_id = str(result[0])  # UUID real
        finally:
            session.close()
        
        
        print(f"📥 Solicitud para crear PPT con {len(graph_ids)} gráficos del usuario {user_id}")

        prs = Presentation()
        
        # Imagen como fondo
        portada = prs.slides.add_slide(prs.slide_layouts[6])
        portada.shapes.add_picture("public/portada.png", 0, 0, width=prs.slide_width, height=prs.slide_height)

        for graph_id in graph_ids:
            chart_path = STORAGE_DIR / user_id / graph_id / "chart.json"
            if not chart_path.exists():
                print(f"⚠️ No se encontró el archivo: {chart_path}")
                continue

            with open(chart_path, "r", encoding="utf-8") as f:
                fig_data = json.load(f)

            fig = pio.from_json(json.dumps(fig_data))
            image_path = f"public/temp/{graph_id}.png"
            fig.write_image(image_path, width=800, height=600)

            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout en blanco
            slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))
            
            # Logo en la esquina superior derecha
            logo_path = "public/logo_oscuro.png"
            logo_width = Inches(1.2)  # ajusta según tamaño real del logo
            logo_top = Inches(0.2)
            logo_left = prs.slide_width - logo_width - Inches(0.2)

            slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width)
        
        # Crear directorio temporal si no existe
        (PUBLIC_DIR / "temp").mkdir(exist_ok=True)

        filename = f"reporte_{uuid.uuid4().hex[:8]}.pptx"
        pptx_path = PUBLIC_DIR / "temp" / filename
        prs.save(pptx_path)

        download_url = f"/public/temp/{filename}"
        print(f"✅ PowerPoint generado: {download_url}")
        return JSONResponse(content={"url": download_url}, status_code=200)

    except Exception as e:
        print(f"❌ Error en create_powerpoint: {e}")
        return JSONResponse(content={"error": str(e)}, status_code=500)
    
 
# Ruta para la página de gráficos
@app.get("/graphs", response_class=HTMLResponse)
async def graphs_page(request: Request, access_token: Optional[str] = Cookie(None)):
    # Si no hay cookie de sesión, manda a login y vuelve a /graphs
    if not access_token:
        return RedirectResponse(url="/login?next=/graphs", status_code=307)
    return templates.TemplateResponse("graphs.html", {"request": request})
    
# Ruta para la página de gráficos
@app.get("/metrics", response_class=HTMLResponse)
async def metrics_page(request: Request):
    return templates.TemplateResponse(
        "metrics.html",
        {"request": request}
    )
    
    

def get_data_analisis(fecha_inicio, fecha_fin):
   
    QUERY = text("""
        SELECT
            t.id AS thread_id,
            t.name AS thread_name,
            u."identifier" AS usuario,
            s.id AS step_id,
            CASE
                WHEN s.type = 'run' AND f.value IS NOT NULL THEN 'feedback'
                ELSE s.type
            END AS type,
            s.name,
            CASE
                WHEN s.type = 'tool' THEN s.input
                ELSE s.output
            END AS texto,
            f.value AS feedback,
            f."comment" AS feedback_comment,
            e."type" AS element_type,
            s."createdAt"
        FROM public.threads t
        JOIN public.steps s ON s."threadId" = t.id
        LEFT JOIN public.users u ON u.id = t."userId"
        LEFT JOIN public.feedbacks f ON f."forId" = s.id
        LEFT JOIN public.elements e ON e."forId" = s.id
        WHERE s."createdAt" BETWEEN :fecha_inicio AND :fecha_fin
        ORDER BY s."createdAt"
    """)
    try:
            session = Session()
            df = pd.read_sql_query(
                QUERY,
                session.bind,
                params={
                    "fecha_inicio": f"{fecha_inicio} 00:00:00",
                    "fecha_fin": f"{fecha_fin} 23:59:59"
                }
            )
            df['feedback'] = df['feedback'].fillna('')
            df = df[~(df['type'] == 'run')]
            df['thread_id'] = df['thread_id'].astype(str)
            df['step_id'] = df['step_id'].astype(str)

            # Conversión sin metadatos:
            registros = df.to_dict(orient="records")
            return registros
    
    except Exception as e:
        raise RuntimeError(f"Error al obtener análisis: {e}")
    finally:
        session.close()
        
        
# Ruta para la página de gráficos
@app.post("/api/analisis", response_class=HTMLResponse)
def analisis_ia(payload: dict = Body(...), user=Depends(get_user_from_cookie)):
    
    fecha_inicio = payload.get("fecha_inicio")
    fecha_fin = payload.get("fecha_fin")

   
    import os
    from dotenv import load_dotenv
    from openai import AzureOpenAI
    
    load_dotenv()
    
    client = AzureOpenAI(
      azure_endpoint=os.environ.get("AZURE_OPENAI_ENDPOINT"),
      api_key=os.environ.get("AZURE_OPENAI_API_KEY"),
      api_version="2024-05-01-preview"
    )
    try:
        registros = get_data_analisis(fecha_inicio, fecha_fin)
        content = json.dumps(registros, ensure_ascii=False)    
          
        completion = client.chat.completions.create(
        model="GPT-4.1",
        temperature=0.5,
        messages=[
            {"role": "system", "content": """
                # Debes analizar un data frame correspondiente a las converzaciones e interacciones de un chat con un Asistente AI
                ## Explicación de los campos a analizar del Data frame:
                
                | Campo            | Comentario                                                                                                                                                    |
                |------------------|---------------------------------------------------------------------------------------------------------------------------------------------------------------|
                | thread_id        | Un thread_id es una conversación, cada conversación tiene varias interacciones que están descritas por su tipo (campo type).                                  |
                | type             | - `user_message`: es la pregunta del usuario<br>- `assistant_message`: es la respuesta del asistente<br>- `tool`: es la herramienta que utilizó el asistente<br>- `feedback`: es un feedback que envía el usuario |
                | texto            | - Si type es `user_message`, se trata de la pregunta del usuario.<br>- Si es `assistant_message`, es la respuesta y razonamiento del asistente.               |
                | usuario          | Identificador del usuario.                                                                                                                                    |
                | feedback         | - Si es `1` es positivo.<br>- Si es `0` es negativo.                                                                                                          |
                | feedback_comment | Es el comentario del usuario al responder el feedback.                                                                                                        |
                | createdAt        | La fecha de creación.                                                                                                                                         |
                
                Basado en un diccionario de consultas que hacen los usuarios, debes analizar lo siguiente:    
                1.- El tenor más frecuente de las consultas del usuario, ¿cuales son las preguntas más fecuentes?
                2.- Los tipos de análisis que solcitan, que información buscan.
                2.- Análisis de sentimiento en general. 
                3.- Intenta contar y calificar las conversaciones 
                
                | **CALIFICACIÓN DE CONVERSACIÓN** | **DESCRIPCIÓN** |
                |-------------------------------|----------------|
                | **Perfecta** | Agente responde sin problemas durante toda la conversación. Genera las consultas a datos adecuadas. |
                | **Aceptable** | Agente responde sin problemas durante toda la conversación. Se genera una o más consultas a los datos sub-óptimas o la búsqueda en la BBDD no genera resultado. *Ejemplos: No se aplica el filtro por consulta, la consulta SQL no devuelve registros en la BBDD.* |
                | **Necesita mejorar** | Agente presenta comportamientos erráticos de forma parcial durante la conversación. No genera una consulta a datos o genera una consulta incorrecta. *Ejemplos: Usuario consulta por algo y el agente busca por otra cosa, el agente no ejecuta la query en algún punto de la interacción.* |
                | **Fallida** | Agente presenta comportamientos erráticos durante toda la conversación. No genera ninguna query frente a las peticiones del usuario o se cae el sistema. *Ejemplos: El agente no ejecuta queries en ningún punto de la interacción, el agente responde en otro idioma, incoherencias.* |


                Tienes libertad para hacer algún tro tu análisis y por supuesto las conclusiones e insights. 
                El formato de respuesta debe estar en HTML, debe quedar listo para incorporarlo entre unas etiquetas <div> existentes,
                por lo que no uses los <div> iniciales, no uses las comillas ``` ni la etiqueta HTML, solo es necesario 
                el string, listo para que el sistema lo inserte.
                
                No hagas una pregunta fional ya que esto se inserta en un reporte fijo.
                """
            },
            {
                "role": "user",
                "content": content
            }  
            ]
        )

        res = completion.choices[0].message.content     
        print(f"res: {res}")   
       
        # Fecha a usar (usa fecha_fin si viene; si no, hoy)
        fecha_reporte = (
            datetime.strptime(fecha_fin, "%Y-%m-%d").date()
            if fecha_fin else date.today()
        )

        html_analisis = (res or "").encode("utf-8", errors="replace").decode("utf-8")

        with engine.begin() as conn:
            conn.execute(
                text("""
                    INSERT INTO public.analisis_ia (fecha, analisis)
                    VALUES (:fecha, :analisis)
                    ON CONFLICT (fecha) DO UPDATE
                    SET analisis = EXCLUDED.analisis,
                        created_at = now();
                """),
                {"fecha": fecha_reporte, "analisis": html_analisis},
            )
        
        return JSONResponse(content={"texto":res},status_code=200)
    
    except Exception as e:
        return JSONResponse(content={"error": str(e)}, status_code=500)
   
    
@app.post("/api/metrics_summary")
def generate_report(payload: dict = Body(...), user=Depends(get_user_from_cookie)):
    cliente = "QAgent"  # o payload.get("cliente")
    fecha_inicio = payload.get("fecha_inicio")
    fecha_fin = payload.get("fecha_fin")

    QUERY = text("""
        SELECT
            t.id AS thread_id,
            t.name AS thread_name,
            u."identifier" AS usuario,
            s.id AS step_id,
            s.start as  inicio,
            s.end as final,
            CASE
                WHEN s.type = 'run' AND f.value IS NOT NULL THEN 'feedback'
                ELSE s.type
            END AS type,
            s.name,
            CASE
                WHEN s.type = 'tool' THEN s.input
                ELSE s.output
            END AS texto,
            f.value AS feedback,
            f."comment" AS feedback_comment,
            e."type" AS element_type,
            s."createdAt"
        FROM public.threads t
        JOIN public.steps s ON s."threadId" = t.id
        LEFT JOIN public.users u ON u.id = t."userId"
        LEFT JOIN public.feedbacks f ON f."forId" = s.id
        LEFT JOIN public.elements e ON e."forId" = s.id
        WHERE s."createdAt" BETWEEN :fecha_inicio AND :fecha_fin
        ORDER BY s."createdAt"
    """)

    try:
        session = Session()
        df = pd.read_sql_query(
            QUERY,
            session.bind,
            params={
                "fecha_inicio": f"{fecha_inicio} 00:00:00",
                "fecha_fin": f"{fecha_fin} 23:59:59"
            }
        )

        df['feedback'] = df['feedback'].fillna('')
        df = df[~(df['type'] == 'run')]

        # ——— Serie temporal de conversaciones por día 
        df['createdAt'] = pd.to_datetime(df['createdAt'])
        df['fecha'] = df['createdAt'].dt.date

        df_ts = df[['thread_id', 'fecha']].drop_duplicates()
        serie = df_ts.groupby('fecha')['thread_id'].nunique().reset_index().sort_values('fecha')

        # Rellenar días sin conversaciones en el rango
        idx = pd.date_range(serie['fecha'].min(), serie['fecha'].max(), freq='D')
        serie = serie.set_index('fecha').reindex(idx, fill_value=0).rename_axis('fecha').reset_index()

        rango_fechas = serie['fecha'].astype(str).tolist()
        rango_conversaciones = serie['thread_id'].tolist()
        # —————————————————————————————————————————————

        # ——— Cálculos generales ya existentes
        feedback_positivo = df[df['feedback'] == 1].shape[0]
        feedback_negativo = df[df['feedback'] == 0].shape[0]

        conversaciones = df['thread_id'].nunique()
        interacciones = df['step_id'].nunique()
        consultas_usuario = df[df['type'] == 'user_message'].shape[0]
        respuestas_asistente = df[df['type'] == 'assistant_message'].shape[0]
        consultas_fuentes = df[df['type'] == 'tool'].shape[0]
        otras_interacciones = interacciones - consultas_usuario - respuestas_asistente - consultas_fuentes

        mix_consultas_usuario = round(consultas_usuario / interacciones, 2)
        mix_respuestas_asistente = round(respuestas_asistente / interacciones, 2)
        mix_consultas_fuentes = round(consultas_fuentes / interacciones, 2)
        mix_otras_interacciones = round(otras_interacciones / interacciones, 2)

        tools = df[df['type'] == 'tool']['name'].value_counts().to_dict()
        elementos_raw = df[~df['element_type'].isna()]['element_type'].value_counts().to_dict()
        elementos = {k.capitalize(): v for k, v in elementos_raw.items()}
        
       
         # ——— Cálculos de tiempo ————————————————————
        df['inicio'] = pd.to_datetime(df['inicio'])
        df['final'] = pd.to_datetime(df['final'])

        # Calcular la duración de cada paso en segundos
        df['duracion'] = (df['final'] - df['inicio']).dt.total_seconds()

        # Agrupar por thread_id y sumar la duración total de cada conversación
        conversaciones_duracion = df.groupby('thread_id')['duracion'].sum()

        # Filtrar para eliminar posibles NaN (conversaciones sin duración)
        conversaciones_duracion = conversaciones_duracion.dropna()

        # Filtrar para dejar solo conversaciones que duran al menos 1 segundo
        conversaciones_duracion = conversaciones_duracion[conversaciones_duracion >= 2]

        # KPI requeridos
        duracion_mas_baja_raw = float(np.round(conversaciones_duracion.min(), 2)) if not conversaciones_duracion.empty else 0        
        duracion_mas_alta_raw = float(np.round(conversaciones_duracion.max(), 2)) if not conversaciones_duracion.empty else 0
        duracion_promedio_raw = float(np.round(conversaciones_duracion.mean(), 2)) if not conversaciones_duracion.empty else 0
        tiempo_total_raw = float(np.round(conversaciones_duracion.sum(), 2)) if not conversaciones_duracion.empty else 0

        
        duracion_mas_baja = formatea_duracion(duracion_mas_baja_raw)
        duracion_mas_alta = formatea_duracion(duracion_mas_alta_raw)
        duracion_promedio = formatea_duracion(duracion_promedio_raw)
        tiempo_total = formatea_duracion(tiempo_total_raw)

        

        return JSONResponse(content={
            "cliente": cliente,
            "fecha_inicio": fecha_inicio,
            "fecha_fin": fecha_fin,
            "conversaciones": conversaciones,
            "interacciones": interacciones,
            "consultas_usuario": consultas_usuario,
            "respuestas_asistente": respuestas_asistente,
            "consultas_fuentes": consultas_fuentes,
            "otras": otras_interacciones,
            "tools": tools,
            "elementos": elementos,
            "mix_consultas_usuario": mix_consultas_usuario,
            "mix_respuestas_asistente": mix_respuestas_asistente,
            "mix_consultas_fuentes": mix_consultas_fuentes,
            "mix_otras_interacciones": mix_otras_interacciones,
            "feedback_positivo": feedback_positivo,
            "feedback_negativo": feedback_negativo,
            "total_registros": len(df),
            "time_series_dates": rango_fechas,
            "time_series_conversaciones": rango_conversaciones,
            "duracion_mas_baja": duracion_mas_baja,
            "duracion_mas_alta": duracion_mas_alta,
            "duracion_promedio": duracion_promedio,
            "tiempo_total":tiempo_total
        }, status_code=200)
        
    
    except Exception as e:
        return JSONResponse(content={"error": str(e)}, status_code=500)
    finally:
        session.close()


@app.post("/api/metrics_today")
def metrics_today(payload: dict = Body(...), user=Depends(get_user_from_cookie)):
    # Fechas clave
    fecha_hoy = date.today()
    fecha_ayer = fecha_hoy - timedelta(days=1)
    mes_inicio = fecha_hoy.replace(day=1)
    print(f"\nmes_inicio: {mes_inicio}\n")

    # 1. Query para HOY
    QUERY_DIA = text("""
        SELECT
            s."threadId" AS thread_id,
            s.id AS step_id,
            s.type,
            s."createdAt"
        FROM public.steps s
        WHERE s."createdAt"::date = :fecha
    """)

    # 2. Query para ACUMULADO del mes (hasta hoy)
    QUERY_MES = text("""
        SELECT
            s."threadId" AS thread_id,
            s.id AS step_id,
            s.type,
            s."createdAt"
        FROM public.steps s
        WHERE s."createdAt"::date BETWEEN :fecha_inicio AND :fecha_fin
    """)

    session = Session()
    try:
        # HOY
        df_hoy = pd.read_sql_query(QUERY_DIA, session.bind, params={"fecha": fecha_hoy})
        # AYER
        df_ayer = pd.read_sql_query(QUERY_DIA, session.bind, params={"fecha": fecha_ayer})
        # ACUMULADO
        df_mes = pd.read_sql_query(QUERY_MES, session.bind, params={"fecha_inicio": mes_inicio, "fecha_fin": fecha_hoy})
        # ACUMULADO hasta AYER
        df_mes_ayer = pd.read_sql_query(QUERY_MES, session.bind, params={"fecha_inicio": mes_inicio, "fecha_fin": fecha_ayer})

        # --- KPIs ---
        def kpis(df):
            conversaciones = df['thread_id'].nunique()
            interacciones = df['step_id'].nunique()
            errores = df[df['type'] == 'error']['step_id'].nunique() if 'error' in df['type'].unique() else 0
            return conversaciones, interacciones, errores

        conv_hoy, int_hoy, err_hoy = kpis(df_hoy)
        conv_ayer, int_ayer, err_ayer = kpis(df_ayer)
        conv_mes, int_mes, _ = kpis(df_mes)
        conv_mes_ayer, int_mes_ayer, _ = kpis(df_mes_ayer)

        # --- Output para ticker ---
        resultado = {
            "hoy": fecha_hoy.strftime("%d-%m-%Y"),
            "conversaciones_hoy": conv_hoy,
            "conversaciones_ayer": conv_ayer,
            "conversaciones_mes": conv_mes,          
            "interacciones_hoy": int_hoy,
            "interacciones_ayer": int_ayer,
            "interacciones_mes": int_mes,
            "interacciones_mes_var": int_mes_ayer,
            "errores_hoy": err_hoy,
        }

        return JSONResponse(content=resultado, status_code=200)

    except Exception as e:
        return JSONResponse(content={"error": str(e)}, status_code=500)
    finally:
        session.close()


@app.post("/api/get_last_report")
def get_last_report(payload: dict = Body(...), user=Depends(get_user_from_cookie)):
    QUERY = text("""
        SELECT resumen, fecha_reporte
        FROM reportes_metrics
        ORDER BY fecha_reporte DESC
        LIMIT 1
    """)

    try:
        session = Session()
        result = session.execute(QUERY).first()
        if result is None:
            return JSONResponse(content={"error": "No hay reportes en la base de datos."}, status_code=404)

        resumen, fecha_reporte = result

        # Si quieres devolver el JSON tal cual fue almacenado:
        import json
        if isinstance(resumen, str):
            resumen = json.loads(resumen)

        resumen["fecha_reporte"] = str(fecha_reporte)

        return JSONResponse(content=resumen, status_code=200)

    except Exception as e:
        return JSONResponse(content={"error": str(e)}, status_code=500)
    finally:
        session.close()

@app.post("/api/get_last_and_previous_report")
def get_last_and_previous_report(payload: dict = Body(...), user=Depends(get_user_from_cookie)):
    import json
    session = Session()
    try:
        # 1) Último
        QUERY_ULTIMO = text("""
            SELECT resumen, fecha_reporte
            FROM reportes_metrics
            ORDER BY fecha_reporte DESC
            LIMIT 1
        """)
        result_ultimo = session.execute(QUERY_ULTIMO).first()

        if not result_ultimo:
            # Estructura estable sin datos
            return JSONResponse(content={"ultimo": None, "anterior": None}, status_code=200)

        resumen_ultimo, fecha_ultimo = result_ultimo
        if isinstance(resumen_ultimo, str):
            try:
                resumen_ultimo = json.loads(resumen_ultimo)
            except Exception:
                resumen_ultimo = {}  # fallback si quedó texto inválido
        if isinstance(resumen_ultimo, dict):
            resumen_ultimo["fecha_reporte"] = str(fecha_ultimo)

        # 2) Anterior
        QUERY_ANTERIOR = text("""
            SELECT resumen, fecha_reporte
            FROM reportes_metrics
            WHERE fecha_reporte < :fecha_ultimo
            ORDER BY fecha_reporte DESC
            LIMIT 1
        """)
        result_anterior = session.execute(QUERY_ANTERIOR, {"fecha_ultimo": fecha_ultimo}).first()

        resumen_anterior = None
        if result_anterior:
            r_anterior, f_anterior = result_anterior
            if isinstance(r_anterior, str):
                try:
                    r_anterior = json.loads(r_anterior)
                except Exception:
                    r_anterior = {}
            if isinstance(r_anterior, dict):
                r_anterior["fecha_reporte"] = str(f_anterior)
            resumen_anterior = r_anterior

        return JSONResponse(content={"ultimo": resumen_ultimo, "anterior": resumen_anterior}, status_code=200)

    except Exception as e:
        return JSONResponse(content={"error": str(e)}, status_code=500)
    finally:
        session.close()

from fastapi import Body, Depends
from fastapi.responses import JSONResponse
from sqlalchemy import text

@app.post("/api/get_last_analisis_ia")
def get_last_analisis_ia(payload: dict = Body(None), user=Depends(get_user_from_cookie)):
    QUERY = text("""
        SELECT id, fecha, analisis
        FROM public.analisis_ia
        ORDER BY fecha DESC, id DESC
        LIMIT 1
    """)
    session = Session()
    try:
        row = session.execute(QUERY).first()
        if row is None:
            return JSONResponse(content={"error": "No hay análisis en la base de datos."}, status_code=404)

        id_, fecha, analisis = row
        # Devolvemos ambos nombres por compatibilidad con el front:
        return JSONResponse(content={
            "id": id_,
            "fecha": str(fecha),
            "analisis": analisis,
            "texto": analisis
        }, status_code=200)

    except Exception as e:
        return JSONResponse(content={"error": str(e)}, status_code=500)
    finally:
        session.close()

@app.post("/api/refresh-graph")
def refresh_graph(payload: dict = Body(...), user=Depends(get_user_from_cookie)):
    print("\n🧪 [refresh-graph] Paso 4 — ejecutar SQL cliente, actualizar chart.json y marcar flags")
    session = None
    try:
        graph_id = payload.get("id")
        if not graph_id:
            print("❌ [refresh-graph] Falta 'id' en el payload")
            return JSONResponse(content={"error": "Falta el parámetro 'id'."}, status_code=400)

        user_identifier = user["identifier"]
        print(f"🆔 [refresh-graph] graph_id = {graph_id}")
        print(f"👤 [refresh-graph] user = {user_identifier}")

        session = Session()
        print("🔌 [refresh-graph] Sesión DB abierta")

        # 1) Traer el element y validar pertenencia
        sql_lookup = text("""
            SELECT
                e.id,
                e."threadId",
                e."objectKey",
                e.pin_sql
            FROM elements e
            JOIN threads t ON e."threadId" = t.id
            WHERE e.id = :graph_id
              AND t."userIdentifier" = :user_identifier
            LIMIT 1
        """)
        row = session.execute(sql_lookup, {"graph_id": graph_id, "user_identifier": user_identifier}).fetchone()

        if not row:
            print("❌ [refresh-graph] No se encontró el gráfico o no pertenece al usuario")
            return JSONResponse(content={"error": "Gráfico no encontrado o no te pertenece."}, status_code=404)

        object_key = row.objectKey
        pin_sql = row.pin_sql

        print(f"📦 [refresh-graph] objectKey = {object_key}")
        print(f"🧾 [refresh-graph] pin_sql {'OK' if pin_sql else 'VACÍO/NULL'}")

        if not object_key:
            print("❌ [refresh-graph] El gráfico no tiene 'objectKey' (ruta del chart.json)")
            return JSONResponse(content={"error": "El gráfico no tiene ruta de chart.json (objectKey)."}, status_code=400)

        if not pin_sql or not str(pin_sql).strip():
            print("ℹ️ [refresh-graph] El gráfico no cuenta con una consulta SQL almacenada (pin_sql vacío)")
            return JSONResponse(
                content={"status": "no_sql", "message": "Tu gráfico no cuenta con una consulta SQL almacenada."},
                status_code=200
            )

        # ---------------------------------------------------------------------
        # 2) Ejecutar la SQL del cliente según .env, mapear a x/y y actualizar chart.json
        # ---------------------------------------------------------------------
        from urllib.parse import quote_plus
        from pathlib import Path
        import os, json
        from datetime import datetime

        def _cfg(k, default=None):
            try:
                return config.get(k, os.getenv(k, default))
            except NameError:
                return os.getenv(k, default)

        tipo = (_cfg("BD_TIPO_UPDATE_GRAPHICS", "MYSQL") or "MYSQL").upper()
        print(f"🔧 [refresh-graph] BD_TIPO_UPDATE_GRAPHICS = {tipo}")

        # --- 2.a Ejecutar consulta en el motor cliente y dejar datos en rows_out/keys_out ---
        rows_out, keys_out = None, None

        if tipo == "MYSQL":
            import mysql.connector
            host = _cfg("DB_HOST")
            port = int(_cfg("DB_MYSQL_PORT", "3306"))
            user = _cfg("DB_USER")
            pwd  = _cfg("DB_PASSWORD")
            name = _cfg("DB_NAME")
            if not all([host, user, name]):
                raise RuntimeError("Credenciales MySQL incompletas (DB_HOST, DB_MYSQL_PORT, DB_USER, DB_PASSWORD, DB_NAME).")

            print(f"🔗 [refresh-graph] Conectando (mysql.connector) a {host}:{port}/{name} …")
            db_conn = mysql.connector.connect(host=host, port=port, user=user, password=pwd, database=name)
            cur = db_conn.cursor(dictionary=True)
            print("▶️  [refresh-graph] Ejecutando pin_sql (MySQL)…")
            cur.execute(pin_sql)
            rows_out = cur.fetchall()   # list[dict]
            keys_out = list(rows_out[0].keys()) if rows_out else []
            print(f"✅ [refresh-graph] pin_sql OK — filas: {len(rows_out)} columnas: {keys_out[:5]}{'…' if len(keys_out)>5 else ''}")
            for i, r in enumerate(rows_out[:3]):
                print(f"   · row[{i}]: {r}")
            cur.close()
            db_conn.close()

        elif tipo in ("MSSQL", "MSSQL_AZURE", "MSSQL_GCP"):
            from sqlalchemy import create_engine
            d_is_azure = tipo in ("MSSQL", "MSSQL_AZURE")
            if d_is_azure:
                server   = _cfg("DB_ASQLS_SERVER")
                database = _cfg("DB_ASQLS_DATABASE")
                user     = _cfg("DB_ASQLS_USERNAME")
                pwd      = _cfg("DB_ASQLS_PASSWORD")
                driver   = _cfg("DB_ASQLS_DRIVER", "ODBC Driver 18 for SQL Server")
                if not all([server, database, user, pwd, driver]):
                    raise RuntimeError("Credenciales Azure SQL incompletas (DB_ASQLS_*).")
            else:
                server   = _cfg("DB_GCP_SQLS_SERVER")
                database = _cfg("DB_CP_SQLS_DATABASE") or _cfg("DB_GCP_SQLS_DATABASE")
                user     = _cfg("DB_DB_GCP_SQLS_USERNAME") or _cfg("DB_GCP_SQLS_USERNAME")
                pwd      = _cfg("DB_GCP_SQLS_PASSWORD")
                driver   = _cfg("DB_GCP_SQLS_DRIVER", "ODBC Driver 18 for SQL Server")
                if not all([server, database, user, pwd, driver]):
                    raise RuntimeError("Credenciales GCP SQL Server incompletas (DB_GCP_SQLS_*).")

            dsn = (
                f"mssql+pyodbc://{quote_plus(user)}:{quote_plus(pwd)}@{quote_plus(server)}/"
                f"{quote_plus(database)}?driver={quote_plus(driver)}&TrustServerCertificate=yes"
            )
            print(f"🔗 [refresh-graph] Conectando a SQL Server {server}/{database} (driver={driver}) …")
            engine_cli = create_engine(dsn, echo=True, pool_pre_ping=True, future=True)

            with engine_cli.connect() as conn:
                print("▶️  [refresh-graph] Ejecutando pin_sql (SQL Server)…")
                result = conn.execute(text(pin_sql))
                rows = result.fetchall()
                cols = result.keys()
                keys_out = list(cols)
                rows_out = [dict(zip(keys_out, row)) for row in rows]
                print(f"✅ [refresh-graph] pin_sql OK — filas: {len(rows_out)} columnas: {keys_out[:5]}{'…' if len(keys_out)>5 else ''}")
                for i, r in enumerate(rows_out[:3]):
                    print(f"   · row[{i}]: {r}")

            try:
                engine_cli.dispose()
            except Exception:
                pass
        else:
            raise RuntimeError(f"BD_TIPO_UPDATE_GRAPHICS='{tipo}' no soportado. Usa MYSQL | MSSQL_AZURE | MSSQL_GCP.")

        # --- 2.b Mapear columnas → x / y ---
        if not rows_out:
            print("ℹ️ [refresh-graph] Consulta devolvió 0 filas; no se actualiza chart.json")
            return JSONResponse(
                content={"status": "no_data", "message": "La consulta no devolvió filas."},
                status_code=200
            )

        lower_cols = [c.lower() for c in keys_out]
        if "mes" in lower_cols and "ventas_mensuales" in lower_cols:
            ix_x = lower_cols.index("mes")
            ix_y = lower_cols.index("ventas_mensuales")
        elif len(keys_out) >= 2:
            ix_x, ix_y = 0, 1
        else:
            return JSONResponse(
                content={"status": "error", "message": "La consulta no posee columnas suficientes para graficar (min 2)."},
                status_code=400
            )

        def _to_num(v):
            try:
                return float(v)
            except Exception:
                return v

        x_values = [list(row.values())[ix_x] for row in rows_out]
        y_values = [_to_num(list(row.values())[ix_y]) for row in rows_out]

        print(f"📊 [refresh-graph] x({len(x_values)})={x_values[:5]}{'…' if len(x_values)>5 else ''}")
        print(f"📈 [refresh-graph] y({len(y_values)})={y_values[:5]}{'…' if len(y_values)>5 else ''}")

        # --- 2.c Abrir chart.json existente ---
        chart_path = (PUBLIC_DIR / "storage" / object_key)
        print(f"🗂️  [refresh-graph] chart_path = {chart_path}")
        if not chart_path.exists():
            return JSONResponse(content={"error": f"No se encontró chart.json en: {str(chart_path)}"}, status_code=400)

        with chart_path.open("r", encoding="utf-8") as f:
            chart = json.load(f)

        if not isinstance(chart, dict) or "data" not in chart or not isinstance(chart["data"], list) or not chart["data"]:
            return JSONResponse(content={"error": "chart.json inválido: no contiene 'data'."}, status_code=400)

        # Actualizar el primer trace
        trace = chart["data"][0]
        trace["x"] = x_values
        trace["y"] = y_values

        # ⚠️ Importante: NO forzamos 'text' — deja que Plotly use texttemplate/format del layout
        if "text" in trace:
            try:
                # Elimina un 'text' viejo para evitar que ensucie el render
                del trace["text"]
            except Exception:
                pass

        # --- 2.d Guardar atómicamente ---
        tmp_path = chart_path.with_suffix(".json.tmp")
        with tmp_path.open("w", encoding="utf-8") as f:
            json.dump(chart, f, ensure_ascii=False)
        tmp_path.replace(chart_path)
        print("💾 [refresh-graph] chart.json actualizado (write tmp → rename)")

        # --- 2.e Marcar update en DB ---
        hoy = datetime.now().strftime("%Y-%m-%d")
        upd = text("""
            UPDATE elements
            SET pin_update = 1,
                pin_update_at = :hoy
            WHERE id = :graph_id
        """)
        session.execute(upd, {"hoy": hoy, "graph_id": graph_id})
        session.commit()
        print(f"🟢 [refresh-graph] elements.pin_update=1, pin_update_at='{hoy}'")

        # Respuesta final
        return JSONResponse(
            content={
                "status": "ok",
                "message": "chart.json actualizado y flags de elemento marcados.",
                "x_len": len(x_values),
                "y_len": len(y_values),
                "pin_update_at": hoy,
                "objectKey": object_key
            },
            status_code=200
        )

    except Exception as e:
        print(f"❌ [refresh-graph] Error inesperado: {e}")
        return JSONResponse(content={"error": str(e)}, status_code=500)
    finally:
        if session is not None:
            session.close()
            print("✅ [refresh-graph] Sesión DB cerrada")




@app.get("/mapa", response_class=HTMLResponse)
def mostrar_mapa(request: Request):
    return templates.TemplateResponse("mapa.html", {"request": request})

# ========== ENDPOINTS PARA REPORTE DE CONVERSACIONES ==========

# 1. Página principal del reporte (solo administradores)
@app.get("/conversations-report", response_class=HTMLResponse)
async def conversations_report_page(request: Request, user = Depends(get_user_from_cookie)):
    """
    Renderiza la página de reporte de conversaciones.
    Solo accesible para usuarios con rol 'admin'.
    """
    verify_admin_role(user)  # Verificar que el usuario sea admin
    return templates.TemplateResponse(
        "conversations_report.html",
        {"request": request}
    )

# 2. API para obtener datos de conversaciones con filtros
@app.get("/api/conversations-data")
def get_conversations_data(
    fecha_inicio: Optional[str] = None,
    fecha_fin: Optional[str] = None,
    usuario: Optional[str] = None,
    user = Depends(get_user_from_cookie)
):
    """
    API que devuelve los datos de conversaciones en formato JSON.
    Soporta filtros por fecha de inicio/fin y usuario específico.
    Los datos incluyen mensajes, herramientas ejecutadas y feedbacks.
    """
    try:
        verify_admin_role(user)
        from QAgent.services.report_service import ReportService
        service = ReportService()

        # Log opcional para debugging
        print(f"📊 Obteniendo conversaciones - Fechas: {fecha_inicio} a {fecha_fin}, Usuario: {usuario}")

        # Cargar threads con todas sus interacciones
        threads = service.load_threads_with_interactions(
            fecha_inicio=fecha_inicio,
            fecha_fin=fecha_fin,
            user_identifier=usuario
        )

        return {"threads": threads, "error": None}
    except Exception as e:
        print(f"❌ Error obteniendo conversaciones: {str(e)}")
        return {"threads": [], "error": str(e)}

# 3. Endpoint de verificación de acceso (para control en frontend)
@app.get("/api/check-reports-access")
def check_reports_access(user = Depends(get_user_from_cookie)):
    """
    Verifica si el usuario actual tiene permisos para acceder a reportes.
    Usado por el frontend para mostrar modal de acceso denegado si no es admin.
    """
    try:
        user_metadata = user.get("metadata", {})
        user_role = user_metadata.get("role", "user").lower()

        if user_role == "admin":
            return {
                "can_access": True,
                "redirect_url": "/conversations-report"
            }
        else:
            return {
                "can_access": False,
                "message": "Acceso Restringido",
                "description": "Esta función está disponible solo para administradores. Contacta al administrador del sistema."
            }
    except Exception:
        return {
            "can_access": False,
            "message": "Error",
            "description": "No se pudo verificar los permisos"
        }
        
# Montar Chainlit en la ruta raíz
mount_chainlit(app=app, target="app.py", path="/") 